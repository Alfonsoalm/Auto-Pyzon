import os
import argparse
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# === Función para crear o añadir a un PowerPoint imágenes de una carpeta con textos debajo ===
def create_presentation_with_images(folder_path, ppt_path=None):
    try:
        if not ppt_path or ppt_path.strip() == "":
            ppt_path = os.path.join(folder_path, "presentacion.pptx")

        if os.path.exists(ppt_path):
            presentation = Presentation(ppt_path)
            if not presentation.slides:
                presentation.slides.add_slide(presentation.slide_layouts[5])  # Añadir diapositiva en blanco
        else:
            presentation = Presentation()
            presentation.slides.add_slide(presentation.slide_layouts[5])  # Añadir diapositiva en blanco

        image_files = sorted(
            [os.path.join(folder_path, f) for f in os.listdir(folder_path) if f.lower().endswith(('png', 'jpg', 'jpeg'))]
        )

        if not image_files:
            return "No se encontraron imagenes en la carpeta."

        image_index = 0
        for slide in presentation.slides:
            if image_index >= len(image_files):
                break

            # Añadir primera imagen (izquierda)
            left = Inches(0.5)
            top = Inches(1.5)
            slide.shapes.add_picture(image_files[image_index], left, top, width=Inches(4))

            # Añadir texto debajo de la primera imagen
            text_box_left = slide.shapes.add_textbox(left, top + Inches(3.5), Inches(4), Inches(0.5))
            text_frame = text_box_left.text_frame
            text_frame.text = "Imagen 1"
            text_frame.paragraphs[0].font.size = Pt(28)
            text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

            image_index += 1

            # Añadir segunda imagen (derecha)
            if image_index < len(image_files):
                left = Inches(5)
                slide.shapes.add_picture(image_files[image_index], left, top, width=Inches(4))

                # Añadir texto debajo de la segunda imagen
                text_box_right = slide.shapes.add_textbox(left, top + Inches(3.5), Inches(4), Inches(0.5))
                text_frame = text_box_right.text_frame
                text_frame.text = "Imagen 2"
                text_frame.paragraphs[0].font.size = Pt(28)
                text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

                image_index += 1

        # Añadir más imágenes a nuevas diapositivas si hay
        while image_index < len(image_files):
            slide = presentation.slides.add_slide(presentation.slide_layouts[5])
            left = Inches(0.5)
            top = Inches(1.5)

            # Añadir primera imagen (izquierda)
            slide.shapes.add_picture(image_files[image_index], left, top, width=Inches(4))

            # Añadir texto debajo de la primera imagen
            text_box_left = slide.shapes.add_textbox(left, top + Inches(3.5), Inches(4), Inches(0.5))
            text_frame = text_box_left.text_frame
            text_frame.text = "Imagen 1"
            text_frame.paragraphs[0].font.size = Pt(28)
            text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

            image_index += 1

            # Añadir segunda imagen (derecha)
            if image_index < len(image_files):
                left = Inches(5)
                slide.shapes.add_picture(image_files[image_index], left, top, width=Inches(4))

                # Añadir texto debajo de la segunda imagen
                text_box_right = slide.shapes.add_textbox(left, top + Inches(3.5), Inches(4), Inches(0.5))
                text_frame = text_box_right.text_frame
                text_frame.text = "Imagen 2"
                text_frame.paragraphs[0].font.size = Pt(28)
                text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

                image_index += 1

        presentation.save(ppt_path)
        return f"Presentacion guardada exitosamente: {ppt_path}"

    except Exception as e:
        return f"Error al procesar la presentacion: {e}"


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Gestor de documentos PowerPoint")
    parser.add_argument("--task", choices=["create_presentation_with_images"], required=True, help="Tarea a realizar.")
    parser.add_argument("--folder", required=True, help="Ruta de la carpeta que contiene las imagenes.")
    parser.add_argument("--output_file", required=False, help="Ruta del archivo PowerPoint de salida.")
    args = parser.parse_args()

    if args.task == "create_presentation_with_images":
        result = create_presentation_with_images(args.folder, args.output_file)
        print(result)
